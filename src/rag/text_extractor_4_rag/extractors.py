# -*- coding: utf-8 -*-
# =============================================================================
# Process Name: Text Extractor for RAG — Core Extraction Engine
# =============================================================================
# Description:
#   TextExtractor class — extracts plain text from 40+ file formats:
#   PDF, DOCX, DOC, XLSX, PPTX, images (OCR), HTML, Markdown, JSON, XML,
#   YAML, archives (ZIP/RAR/7Z/TAR), source code, EML, MSG, EPUB, ODT, RTF.
#   Also supports web page extraction with optional JS rendering.
#
# File: src/rag/text_extractor_4_rag/extractors.py
# Project: FastApiFoundry (Docker)
# Version: 0.6.0
# Author: hypo69
# Copyright: © 2026 hypo69
# =============================================================================

import asyncio
import concurrent.futures
import io
import logging
import os
import shutil
import subprocess
import tarfile
import tempfile
import threading
import time
import xml.etree.ElementTree as ET
import zipfile
from pathlib import Path
from typing import Any, Dict, List, Optional

# Импорты для архивов
try:
    import rarfile
except ImportError:
    rarfile = None

try:
    import py7zr
except ImportError:
    py7zr = None

# Импорты для различных форматов
try:
    import pdfplumber
    import PyPDF2
except ImportError:
    PyPDF2 = None
    pdfplumber = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    import pandas as pd
except ImportError:
    pd = None

try:
    import pytesseract
    from PIL import Image
except ImportError:
    Image = None
    pytesseract = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from bs4 import BeautifulSoup
except ImportError:
    BeautifulSoup = None

try:
    import markdown
except ImportError:
    markdown = None

try:
    from odf.opendocument import load
    from odf.teletype import extractText
    from odf.text import P
except ImportError:
    load = None
    P = None
    extractText = None

try:
    from striprtf.striprtf import rtf_to_text
except ImportError:
    rtf_to_text = None

try:
    import yaml
except ImportError:
    yaml = None

# Веб-экстракция (новое в v1.10.0)
try:
    import ipaddress
    from urllib.parse import urljoin, urlparse

    import requests
except ImportError:
    requests = None
    urljoin = None
    urlparse = None
    ipaddress = None

# Playwright для JS-рендеринга (новое в v1.10.1)
try:
    from playwright.sync_api import sync_playwright
except ImportError:
    sync_playwright = None

from .config import settings
from .utils import get_file_extension, is_archive_format, is_supported_format

logger = logging.getLogger(__name__)


class TextExtractor:
    """Класс для извлечения текста из файлов различных форматов."""

    def __init__(self):
        """Инициализация экстрактора текста."""
        self.ocr_languages = settings.OCR_LANGUAGES
        self.timeout = settings.PROCESSING_TIMEOUT_SECONDS
        # Apply explicit tesseract path from config.json → text_extractor.tesseract_cmd.
        # Set automatically by install\install-tesseract.ps1.
        if pytesseract:
            tess_cmd = settings.TESSERACT_CMD
            if tess_cmd and os.path.isfile(tess_cmd):
                pytesseract.pytesseract.tesseract_cmd = tess_cmd
        self._thread_pool = concurrent.futures.ThreadPoolExecutor(max_workers=4)

    def extract_text(self, file_content: bytes, filename: str) -> List[Dict[str, Any]]:
        """Основной метод извлечения текста (теперь синхронный для выполнения в threadpool)."""
        # Проверка, является ли файл архивом
        if is_archive_format(filename, settings.SUPPORTED_FORMATS):
            return self._extract_from_archive(file_content, filename)

        # Проверка поддержки формата
        if not is_supported_format(filename, settings.SUPPORTED_FORMATS):
            raise ValueError(f"Unsupported file format: {filename}")

        # Проверка MIME-типа для безопасности (синхронная операция)
        is_valid_mime = self._check_mime_type(file_content, filename)

        if not is_valid_mime:
            logger.warning(f"MIME-тип файла {filename} не соответствует расширению")
            # Не блокируем, но предупреждаем

        extension = get_file_extension(filename)

        # Проверка, что extension не None
        if not extension:
            raise ValueError(f"Could not determine file extension for: {filename}")

        try:
            # Извлечение текста синхронно
            text = self._extract_text_by_format(file_content, extension, filename)

            # Возвращаем массив с одним элементом для единообразия
            return [
                {
                    "filename": filename,
                    "path": filename,
                    "size": len(file_content),
                    "type": extension,
                    "text": text.strip() if text else "",
                }
            ]

        except Exception as e:
            logger.error(f"Ошибка при извлечении текста из {filename}: {str(e)}")
            raise ValueError(f"Error extracting text: {str(e)}")

    def _extract_text_by_format(
        self, content: bytes, extension: str, filename: str
    ) -> str:
        """Извлечение текста в зависимости от формата (синхронная версия)."""
        # Создаем словарь сопоставления расширений с методами извлечения
        extraction_methods = self._get_extraction_methods_mapping()

        # Проверяем, является ли файл исходным кодом
        source_code_extensions = settings.SUPPORTED_FORMATS.get("source_code", [])
        if extension in source_code_extensions:
            return self._extract_from_source_code_sync(content, extension, filename)

        # Ищем подходящий метод извлечения
        extractor_method = extraction_methods.get(extension)
        if extractor_method:
            return extractor_method(content)

        # Проверяем группы расширений
        for extensions_group, method in self._get_group_extraction_methods():
            if extension in extensions_group:
                return method(content)

        raise ValueError(f"Unsupported file format: {extension}")

    def _get_extraction_methods_mapping(self) -> dict:
        """Получение словаря сопоставления расширений с методами извлечения."""
        return {
            "pdf": self._extract_from_pdf_sync,
            "docx": self._extract_from_docx_sync,
            "doc": self._extract_from_doc_sync,
            "csv": self._extract_from_csv_sync,
            "pptx": self._extract_from_pptx_sync,
            "ppt": self._extract_from_ppt_sync,
            "txt": self._extract_from_txt_sync,
            "json": self._extract_from_json_sync,
            "rtf": self._extract_from_rtf_sync,
            "odt": self._extract_from_odt_sync,
            "xml": self._extract_from_xml_sync,
            "epub": self._extract_from_epub_sync,
            "eml": self._extract_from_eml_sync,
            "msg": self._extract_from_msg_sync,
        }

    def _get_group_extraction_methods(self) -> list:
        """Получение списка групп расширений с соответствующими методами."""
        return [
            (["xls", "xlsx"], self._extract_from_excel_sync),
            (
                ["jpg", "jpeg", "png", "tiff", "tif", "bmp", "gif"],
                self._extract_from_image_sync,
            ),
            (["html", "htm"], self._extract_from_html_sync),
            (["md", "markdown"], self._extract_from_markdown_sync),
            (["yaml", "yml"], self._extract_from_yaml_sync),
        ]

    def _extract_from_pdf_sync(self, content: bytes) -> str:
        """Синхронное извлечение текста из PDF."""
        if not pdfplumber:
            raise ImportError("pdfplumber не установлен")

        text_parts = []
        temp_file_path = None

        try:
            with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as temp_file:
                temp_file.write(content)
                temp_file_path = temp_file.name

            with pdfplumber.open(temp_file_path) as pdf:
                for page_num, page in enumerate(pdf.pages, 1):
                    page_texts = self._extract_pdf_page_content(page, page_num)
                    text_parts.extend(page_texts)

            return "\n\n".join(text_parts)

        except Exception as e:
            logger.error(f"Ошибка при обработке PDF: {str(e)}")
            raise ValueError(f"Error processing PDF: {str(e)}")
        finally:
            self._cleanup_temp_file(temp_file_path)

    def _extract_pdf_page_content(self, page, page_num: int) -> list:
        """Извлечение содержимого страницы PDF."""
        page_texts = []

        # Извлечение текста со страницы
        page_text = page.extract_text()
        if page_text:
            page_texts.append(f"[Страница {page_num}]\n{page_text}")

        # Извлечение изображений и OCR
        if page.images:
            page_texts.extend(self._extract_pdf_page_images(page))

        return page_texts

    def _extract_pdf_page_images(self, page) -> list:
        """Извлечение текста из изображений на странице PDF."""
        image_texts = []

        for img_idx, img in enumerate(page.images):
            try:
                image_text = self._ocr_from_pdf_image_sync(page, img)
                if image_text.strip():
                    image_texts.append(f"[Изображение {img_idx + 1}]\n{image_text}")
            except Exception as e:
                logger.warning(f"Ошибка OCR изображения {img_idx + 1}: {str(e)}")

        return image_texts

    def _cleanup_temp_file(self, temp_file_path: str) -> None:
        """Безопасное удаление временного файла."""
        if temp_file_path and os.path.exists(temp_file_path):
            try:
                os.unlink(temp_file_path)
            except OSError as e:
                logger.warning(
                    f"Не удалось удалить временный файл {temp_file_path}: {str(e)}"
                )

    def _extract_from_docx_sync(self, content: bytes) -> str:
        """Синхронное извлечение текста из DOCX с полным извлечением согласно п.3.3 ТЗ."""
        if not Document:
            raise ImportError("python-docx не установлен")

        try:
            doc = Document(io.BytesIO(content))
            text_parts = []

            # Извлекаем различные части документа
            text_parts.extend(self._extract_docx_paragraphs(doc))
            text_parts.extend(self._extract_docx_tables(doc))
            text_parts.extend(self._extract_docx_headers_footers(doc))
            text_parts.extend(self._extract_docx_footnotes(doc))
            text_parts.extend(self._extract_docx_comments(doc))

            return "\n\n".join(text_parts)

        except Exception as e:
            logger.error(f"Ошибка при обработке DOCX: {str(e)}")
            raise ValueError(f"Error processing DOCX: {str(e)}")

    def _extract_docx_paragraphs(self, doc) -> list:
        """Извлечение основного текста из параграфов DOCX."""
        text_parts = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_parts.append(paragraph.text)
        return text_parts

    def _extract_docx_tables(self, doc) -> list:
        """Извлечение текста из таблиц DOCX."""
        text_parts = []
        for table in doc.tables:
            table_text = []
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    row_text.append(cell.text.strip())
                table_text.append("\t".join(row_text))

            if table_text:
                text_parts.append("\n".join(table_text))
        return text_parts

    def _extract_docx_headers_footers(self, doc) -> list:
        """Извлечение текста из колонтитулов DOCX."""
        text_parts = []
        for section in doc.sections:
            # Извлечение заголовка
            if section.header:
                header_text = self._extract_section_text(section.header.paragraphs)
                if header_text:
                    text_parts.append(
                        f"[Колонтитул - Заголовок]\n{' '.join(header_text)}"
                    )

            # Извлечение подвала
            if section.footer:
                footer_text = self._extract_section_text(section.footer.paragraphs)
                if footer_text:
                    text_parts.append(f"[Колонтитул - Подвал]\n{' '.join(footer_text)}")
        return text_parts

    def _extract_section_text(self, paragraphs) -> list:
        """Извлечение текста из параграфов секции."""
        text_parts = []
        for paragraph in paragraphs:
            if paragraph.text.strip():
                text_parts.append(paragraph.text)
        return text_parts

    def _extract_docx_footnotes(self, doc) -> list:
        """Извлечение сносок из DOCX."""
        text_parts = []
        try:
            if hasattr(doc, "footnotes") and doc.footnotes:
                footnotes_text = []
                for footnote in doc.footnotes:
                    if hasattr(footnote, "paragraphs"):
                        footnote_text = self._extract_section_text(footnote.paragraphs)
                        footnotes_text.extend(footnote_text)
                if footnotes_text:
                    text_parts.append(f"[Сноски]\n{' '.join(footnotes_text)}")
        except Exception as e:
            logger.debug(f"Не удалось извлечь сноски из DOCX: {str(e)}")
        return text_parts

    def _extract_docx_comments(self, doc) -> list:
        """Извлечение комментариев из DOCX."""
        text_parts = []
        try:
            if hasattr(doc, "comments") and doc.comments:
                comments_text = []
                for comment in doc.comments:
                    if hasattr(comment, "paragraphs"):
                        comment_text = self._extract_section_text(comment.paragraphs)
                        comments_text.extend(comment_text)
                if comments_text:
                    text_parts.append(f"[Комментарии]\n{' '.join(comments_text)}")
        except Exception as e:
            logger.debug(f"Не удалось извлечь комментарии из DOCX: {str(e)}")
        return text_parts

    def _extract_from_doc_sync(self, content: bytes) -> str:
        """Синхронное извлечение текста из DOC через конвертацию в DOCX с помощью LibreOffice."""
        if not Document:
            raise ImportError("python-docx не установлен")

        try:
            # Создаем временные файлы
            with tempfile.NamedTemporaryFile(suffix=".doc", delete=False) as temp_doc:
                temp_doc.write(content)
                temp_doc_path = temp_doc.name

            # Создаем временную директорию для вывода
            temp_dir = tempfile.mkdtemp()

            try:
                # Конвертируем .doc в .docx с помощью LibreOffice с ограничениями ресурсов
                result = run_subprocess_with_limits(
                    command=[
                        "libreoffice",
                        "--headless",
                        "--convert-to",
                        "docx",
                        "--outdir",
                        temp_dir,
                        temp_doc_path,
                    ],
                    timeout=30,
                    memory_limit=settings.MAX_LIBREOFFICE_MEMORY,
                    capture_output=True,
                    text=True,
                )

                if result.returncode != 0:
                    logger.error(f"LibreOffice conversion failed: {result.stderr}")
                    raise ValueError("Failed to convert DOC to DOCX")

                # Находим сконвертированный файл
                doc_filename = os.path.splitext(os.path.basename(temp_doc_path))[0]
                docx_path = os.path.join(temp_dir, f"{doc_filename}.docx")

                if not os.path.exists(docx_path):
                    raise ValueError("Converted DOCX file not found")

                # Читаем сконвертированный DOCX файл
                with open(docx_path, "rb") as docx_file:
                    docx_content = docx_file.read()

                # Используем синхронный метод для извлечения текста из DOCX
                text = self._extract_from_docx_sync(docx_content)

                return text

            finally:
                # Очищаем временные файлы
                try:
                    os.unlink(temp_doc_path)
                except Exception as e:
                    logger.warning(
                        f"Не удалось удалить временный файл {temp_doc_path}: {e}"
                    )

                try:
                    import shutil

                    shutil.rmtree(temp_dir, ignore_errors=True)
                except Exception as e:
                    logger.warning(
                        f"Не удалось удалить временную директорию {temp_dir}: {e}"
                    )

        except subprocess.TimeoutExpired:
            logger.error("LibreOffice conversion timeout")
            raise ValueError("DOC conversion timeout")
        except MemoryError as e:
            logger.error(f"LibreOffice превысил лимит памяти: {str(e)}")
            raise ValueError("DOC conversion failed: memory limit exceeded")
        except Exception as e:
            logger.error(f"Ошибка при обработке DOC: {str(e)}")
            raise ValueError(f"Error processing DOC: {str(e)}")

    def _extract_from_excel_sync(self, content: bytes) -> str:
        """Синхронное извлечение данных из Excel файлов."""
        if not pd:
            raise ImportError("pandas не установлен")

        try:
            excel_data = pd.read_excel(io.BytesIO(content), sheet_name=None)
            text_parts = []

            for sheet_name, df in excel_data.items():
                text_parts.append(f"[Лист: {sheet_name}]")
                text_parts.append(df.to_csv(index=False))

            return "\n\n".join(text_parts)

        except Exception as e:
            logger.error(f"Ошибка при обработке Excel: {str(e)}")
            raise ValueError(f"Error processing Excel: {str(e)}")

    def _extract_from_csv_sync(self, content: bytes) -> str:
        """Синхронное извлечение данных из CSV файлов."""
        if not pd:
            raise ImportError("pandas не установлен")

        try:
            df = pd.read_csv(io.BytesIO(content))
            return df.to_csv(index=False)

        except Exception as e:
            logger.error(f"Ошибка при обработке CSV: {str(e)}")
            raise ValueError(f"Error processing CSV: {str(e)}")

    def _extract_from_pptx_sync(self, content: bytes) -> str:
        """Синхронное извлечение текста из PPTX с полным извлечением согласно п.3.3 ТЗ."""
        if not Presentation:
            raise ImportError("python-pptx не установлен")

        try:
            prs = Presentation(io.BytesIO(content))
            text_parts = []

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                slide_text.append(f"[Слайд {slide_num}]")

                # Извлечение текста из фигур слайда
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)

                # Извлечение заметок спикера - согласно п.3.3 ТЗ
                try:
                    if hasattr(slide, "notes_slide") and slide.notes_slide:
                        notes_text = []
                        # Извлечение заметок из текстовых фигур
                        for shape in slide.notes_slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                # Фильтруем стандартные заголовки PowerPoint
                                if shape.text.strip() not in ["Заметки", "Notes"]:
                                    notes_text.append(shape.text.strip())

                        if notes_text:
                            slide_text.append(
                                f"[Заметки спикера]\n{' '.join(notes_text)}"
                            )
                except Exception as e:
                    logger.debug(
                        f"Не удалось извлечь заметки спикера со слайда {slide_num}: {str(e)}"
                    )

                if len(slide_text) > 1:  # Больше чем просто заголовок слайда
                    text_parts.append("\n".join(slide_text))

            return "\n\n".join(text_parts)

        except Exception as e:
            logger.error(f"Ошибка при обработке PPTX: {str(e)}")
            raise ValueError(f"Error processing PPTX: {str(e)}")

    def _extract_from_ppt_sync(self, content: bytes) -> str:
        """Синхронное извлечение текста из PPT через конвертацию в PPTX с помощью LibreOffice."""
        if not Presentation:
            raise ImportError("python-pptx не установлен")

        try:
            # Создаем временные файлы
            with tempfile.NamedTemporaryFile(suffix=".ppt", delete=False) as temp_ppt:
                temp_ppt.write(content)
                temp_ppt_path = temp_ppt.name

            # Создаем временную директорию для вывода
            temp_dir = tempfile.mkdtemp()

            try:
                # Конвертируем .ppt в .pptx с помощью LibreOffice с ограничениями ресурсов
                result = run_subprocess_with_limits(
                    command=[
                        "libreoffice",
                        "--headless",
                        "--convert-to",
                        "pptx",
                        "--outdir",
                        temp_dir,
                        temp_ppt_path,
                    ],
                    timeout=30,
                    memory_limit=settings.MAX_LIBREOFFICE_MEMORY,
                    capture_output=True,
                    text=True,
                )

                if result.returncode != 0:
                    logger.error(f"LibreOffice conversion failed: {result.stderr}")
                    raise ValueError("Failed to convert PPT to PPTX")

                # Находим сконвертированный файл
                ppt_filename = os.path.splitext(os.path.basename(temp_ppt_path))[0]
                pptx_path = os.path.join(temp_dir, f"{ppt_filename}.pptx")

                if not os.path.exists(pptx_path):
                    raise ValueError("Converted PPTX file not found")

                # Читаем сконвертированный PPTX файл
                with open(pptx_path, "rb") as pptx_file:
                    pptx_content = pptx_file.read()

                # Используем синхронный метод для извлечения текста из PPTX
                text = self._extract_from_pptx_sync(pptx_content)

                return text

            finally:
                # Очищаем временные файлы
                try:
                    os.unlink(temp_ppt_path)
                except Exception as e:
                    logger.warning(
                        f"Не удалось удалить временный файл {temp_ppt_path}: {e}"
                    )

                try:
                    import shutil

                    shutil.rmtree(temp_dir, ignore_errors=True)
                except Exception as e:
                    logger.warning(
                        f"Не удалось удалить временную директорию {temp_dir}: {e}"
                    )

        except subprocess.TimeoutExpired:
            logger.error("LibreOffice conversion timeout")
            raise ValueError("PPT conversion timeout")
        except MemoryError as e:
            logger.error(f"LibreOffice превысил лимит памяти: {str(e)}")
            raise ValueError("PPT conversion failed: memory limit exceeded")
        except Exception as e:
            logger.error(f"Ошибка при обработке PPT: {str(e)}")
            raise ValueError(f"Error processing PPT: {str(e)}")

    def _extract_from_txt_sync(self, content: bytes) -> str:
        """Синхронное извлечение текста из TXT файлов."""
        try:
            return self._decode_text_content(content)
        except Exception as e:
            logger.error(f"Ошибка при обработке TXT: {str(e)}")
            raise ValueError(f"Error processing TXT: {str(e)}")

    def _decode_text_content(self, content: bytes) -> str:
        """Декодирование содержимого с автоопределением кодировки."""
        encodings = self._get_encoding_list()

        for encoding in encodings:
            decoded_text = self._try_decode_with_encoding(content, encoding)
            if decoded_text is not None:
                return decoded_text

        # Если не удалось декодировать ни одной кодировкой, используем замещение символов
        logger.warning(
            "Не удалось определить кодировку файла, используем UTF-8 с заменой символов"
        )
        return content.decode("utf-8", errors="replace")

    def _get_encoding_list(self) -> list:
        """Получение списка кодировок для проверки."""
        return [
            "utf-8",  # Стандартная UTF-8
            "mac-cyrillic",  # Macintosh кодировка для кириллицы
            "cp1251",  # Windows-1251 (основная кодировка Windows для русского)
            "windows-1251",  # Альтернативное название для cp1251
            "koi8-r",  # КОИ-8 (старая советская кодировка)
            "cp866",  # DOS кодировка для русского
            "iso-8859-5",  # ISO кодировка для кириллицы
            "utf-16",  # UTF-16 (иногда используется в Windows)
            "utf-16le",  # UTF-16 Little Endian
            "utf-16be",  # UTF-16 Big Endian
            "latin-1",  # ISO-8859-1 (fallback)
            "ascii",  # ASCII (базовая кодировка)
        ]

    def _try_decode_with_encoding(self, content: bytes, encoding: str) -> str:
        """Попытка декодирования с проверкой качества."""
        try:
            decoded_text = content.decode(encoding)

            if not self._is_decoding_quality_good(decoded_text):
                return None

            if not self._is_mac_cyrillic_valid(decoded_text, encoding):
                return None

            return decoded_text
        except UnicodeError:
            return None

    def _is_decoding_quality_good(self, text: str) -> bool:
        """Проверка качества декодирования по количеству заменяющих символов."""
        if "�" in text:
            replacement_ratio = text.count("�") / len(text)
            return replacement_ratio <= 0.1  # Не больше 10% заменяющих символов
        return True

    def _is_mac_cyrillic_valid(self, text: str, encoding: str) -> bool:
        """Дополнительная валидация для mac-cyrillic кодировки."""
        if encoding != "mac-cyrillic" or not text:
            return True

        if self._has_suspicious_start_chars(text):
            return False

        return self._has_valid_cyrillic_ratio(text)

    def _has_suspicious_start_chars(self, text: str) -> bool:
        """Проверка на подозрительные символы в начале текста."""
        suspicious_chars = [
            '"',
            "'",
            "`",
            "«",
            "»",
            '"',
            '"',
            """, """,
            chr(8220),
            chr(8221),
        ]
        return len(text) > 1 and text[0] in suspicious_chars

    def _has_valid_cyrillic_ratio(self, text: str) -> bool:
        """Проверка соотношения кириллицы и латиницы."""
        cyrillic_count = sum(1 for char in text if "\u0400" <= char <= "\u04ff")
        latin_count = sum(1 for char in text if "a" <= char.lower() <= "z")
        total_letters = cyrillic_count + latin_count

        if total_letters == 0:
            return True

        # Если кириллица составляет менее 70% и она есть, то это подозрительно
        return not (cyrillic_count / total_letters < 0.7 and cyrillic_count > 0)

    def _extract_from_source_code_sync(
        self, content: bytes, extension: str, filename: str
    ) -> str:
        """Синхронное извлечение текста из файлов исходного кода."""
        try:
            # Декодируем содержимое файла
            text = self._decode_text_content(content)

            # Получаем информацию о языке программирования и форматируем результат
            return self._format_source_code_output(text, extension, filename)

        except Exception as e:
            logger.error(f"Ошибка при обработке исходного кода {filename}: {str(e)}")
            raise ValueError(f"Error processing source code: {str(e)}")

    def _format_source_code_output(
        self, text: str, extension: str, filename: str
    ) -> str:
        """Форматирование вывода для файлов исходного кода."""
        language = self._get_programming_language(extension)
        header = self._create_source_code_header(language, filename, text)
        return header + "=" * 50 + "\n" + text

    def _get_programming_language(self, extension: str) -> str:
        """Определение языка программирования по расширению файла."""
        language_map = self._get_language_map()
        return language_map.get(extension.lower(), "Source Code")

    def _get_language_map(self) -> dict:
        """Получение словаря соответствия расширений языкам программирования."""
        return {
            # Python
            "py": "Python",
            "pyx": "Python",
            "pyi": "Python",
            "pyw": "Python",
            # JavaScript/TypeScript
            "js": "JavaScript",
            "jsx": "JavaScript",
            "ts": "TypeScript",
            "tsx": "TypeScript",
            "mjs": "JavaScript",
            "cjs": "JavaScript",
            # Java
            "java": "Java",
            "jav": "Java",
            # C/C++
            "c": "C",
            "cpp": "C++",
            "cxx": "C++",
            "cc": "C++",
            "c++": "C++",
            "h": "C Header",
            "hpp": "C++ Header",
            "hxx": "C++ Header",
            "h++": "C++ Header",
            # C#
            "cs": "C#",
            "csx": "C#",
            # PHP
            "php": "PHP",
            "php3": "PHP",
            "php4": "PHP",
            "php5": "PHP",
            "phtml": "PHP",
            # Ruby
            "rb": "Ruby",
            "rbw": "Ruby",
            "rake": "Ruby",
            "gemspec": "Ruby",
            # Go
            "go": "Go",
            "mod": "Go Module",
            "sum": "Go Sum",
            # Rust
            "rs": "Rust",
            "rlib": "Rust Library",
            # Swift
            "swift": "Swift",
            # Kotlin
            "kt": "Kotlin",
            "kts": "Kotlin Script",
            # Scala
            "scala": "Scala",
            "sc": "Scala",
            # R
            "r": "R",
            "R": "R",
            "rmd": "R Markdown",
            "Rmd": "R Markdown",
            # SQL
            "sql": "SQL",
            "ddl": "SQL DDL",
            "dml": "SQL DML",
            # Shell
            "sh": "Shell",
            "bash": "Bash",
            "zsh": "Zsh",
            "fish": "Fish",
            "ksh": "Ksh",
            "csh": "Csh",
            "tcsh": "Tcsh",
            # PowerShell
            "ps1": "PowerShell",
            "psm1": "PowerShell Module",
            "psd1": "PowerShell Data",
            # Perl
            "pl": "Perl",
            "pm": "Perl Module",
            "pod": "Perl Documentation",
            "t": "Perl Test",
            # Lua
            "lua": "Lua",
            # 1C and OneScript
            "bsl": "1C:Enterprise",
            "os": "OneScript",
            # Configuration files
            "ini": "INI Config",
            "cfg": "Config",
            "conf": "Config",
            "config": "Config",
            "toml": "TOML",
            "properties": "Properties",
            # Web
            "css": "CSS",
            "scss": "SCSS",
            "sass": "Sass",
            "less": "Less",
            "styl": "Stylus",
            # Markup
            "tex": "LaTeX",
            "latex": "LaTeX",
            "rst": "reStructuredText",
            "adoc": "AsciiDoc",
            "asciidoc": "AsciiDoc",
            # Data
            "jsonl": "JSON Lines",
            "ndjson": "NDJSON",
            "jsonc": "JSON with Comments",
            # Docker
            "dockerfile": "Dockerfile",
            "containerfile": "Containerfile",
            # Makefile
            "makefile": "Makefile",
            "mk": "Makefile",
            "mak": "Makefile",
            # Git
            "gitignore": "Git Ignore",
            "gitattributes": "Git Attributes",
            "gitmodules": "Git Modules",
        }

    def _create_source_code_header(
        self, language: str, filename: str, text: str
    ) -> str:
        """Создание заголовка для файла исходного кода."""
        header = f"=== {language} File: {filename} ===\n"

        lines = text.split("\n")
        line_count = len(lines)
        header += f"Lines: {line_count}\n"

        # Если файл слишком длинный, добавляем предупреждение
        if line_count > 1000:
            header += f"Warning: Large file with {line_count} lines\n"

        return header

    def _extract_from_html_sync(self, content: bytes) -> str:
        """Синхронное извлечение текста из HTML."""
        if not BeautifulSoup:
            raise ImportError("beautifulsoup4 не установлен")

        try:
            text = content.decode("utf-8", errors="replace")
            soup = BeautifulSoup(text, "html.parser")

            # Удаление script и style тегов
            for script in soup(["script", "style"]):
                script.decompose()

            # Получение текста
            text = soup.get_text()

            # Очистка от лишних пробелов
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            return "\n".join(chunk for chunk in chunks if chunk)

        except Exception as e:
            logger.error(f"Ошибка при обработке HTML: {str(e)}")
            raise ValueError(f"Error processing HTML: {str(e)}")

    def _extract_from_markdown_sync(self, content: bytes) -> str:
        """Синхронное извлечение текста из Markdown."""
        try:
            text = content.decode("utf-8", errors="replace")

            if markdown:
                # Конвертация в HTML и извлечение текста
                html = markdown.markdown(text)
                if BeautifulSoup:
                    soup = BeautifulSoup(html, "html.parser")
                    return soup.get_text()

            # Если markdown не установлен, возвращаем как есть
            return text

        except Exception as e:
            logger.error(f"Ошибка при обработке Markdown: {str(e)}")
            raise ValueError(f"Error processing Markdown: {str(e)}")

    def _extract_from_json_sync(self, content: bytes) -> str:
        """Синхронное извлечение текста из JSON."""
        import json

        try:
            text = content.decode("utf-8", errors="replace")
            data = json.loads(text)

            # Рекурсивное извлечение всех строковых значений
            def extract_strings(obj, path=""):
                strings = []
                if isinstance(obj, dict):
                    for key, value in obj.items():
                        new_path = f"{path}.{key}" if path else key
                        strings.extend(extract_strings(value, new_path))
                elif isinstance(obj, list):
                    for i, value in enumerate(obj):
                        new_path = f"{path}[{i}]" if path else f"[{i}]"
                        strings.extend(extract_strings(value, new_path))
                elif isinstance(obj, str):
                    if obj.strip():
                        strings.append(f"{path}: {obj}")
                return strings

            strings = extract_strings(data)
            return "\n".join(strings)

        except Exception as e:
            logger.error(f"Ошибка при обработке JSON: {str(e)}")
            raise ValueError(f"Error processing JSON: {str(e)}")

    def _extract_from_rtf_sync(self, content: bytes) -> str:
        """Синхронное извлечение текста из RTF."""
        if not rtf_to_text:
            raise ImportError("striprtf не установлен")

        try:
            text = content.decode("utf-8", errors="replace")
            plain_text = rtf_to_text(text)
            return plain_text

        except Exception as e:
            logger.error(f"Ошибка при обработке RTF: {str(e)}")
            raise ValueError(f"Error processing RTF: {str(e)}")

    def _extract_from_xml_sync(self, content: bytes) -> str:
        """Синхронное извлечение текста из XML."""
        try:
            text = content.decode("utf-8", errors="replace")
            root = ET.fromstring(text)

            # Рекурсивное извлечение всех элементов и атрибутов
            def extract_from_element(elem, path=""):
                strings = []

                current_path = f"{path}.{elem.tag}" if path else elem.tag

                # Добавляем текст элемента
                if elem.text and elem.text.strip():
                    strings.append(f"{current_path}: {elem.text.strip()}")

                # Добавляем атрибуты
                for attr_name, attr_value in elem.attrib.items():
                    if attr_value.strip():
                        strings.append(f"{current_path}@{attr_name}: {attr_value}")

                # Рекурсивно обрабатываем дочерние элементы
                for child in elem:
                    strings.extend(extract_from_element(child, current_path))

                return strings

            strings = extract_from_element(root)
            return "\n".join(strings)

        except Exception as e:
            logger.error(f"Ошибка при обработке XML: {str(e)}")
            raise ValueError(f"Error processing XML: {str(e)}")

    def _extract_from_yaml_sync(self, content: bytes) -> str:
        """Синхронное извлечение текста из YAML."""
        if not yaml:
            raise ImportError("PyYAML не установлен")

        try:
            text = content.decode("utf-8", errors="replace")
            data = yaml.safe_load(text)
            strings = self._extract_yaml_strings(data)
            return "\n".join(strings)

        except Exception as e:
            logger.error(f"Ошибка при обработке YAML: {str(e)}")
            raise ValueError(f"Error processing YAML: {str(e)}")

    def _extract_yaml_strings(self, obj, path="") -> list:
        """Рекурсивное извлечение всех строковых значений из YAML."""
        strings = []

        if isinstance(obj, dict):
            strings.extend(self._extract_yaml_dict_strings(obj, path))
        elif isinstance(obj, list):
            strings.extend(self._extract_yaml_list_strings(obj, path))
        elif isinstance(obj, str) and obj.strip():
            strings.append(f"{path}: {obj}")

        return strings

    def _extract_yaml_dict_strings(self, obj_dict: dict, path: str) -> list:
        """Извлечение строк из словаря YAML."""
        strings = []
        for key, value in obj_dict.items():
            new_path = f"{path}.{key}" if path else key
            strings.extend(self._extract_yaml_strings(value, new_path))
        return strings

    def _extract_yaml_list_strings(self, obj_list: list, path: str) -> list:
        """Извлечение строк из списка YAML."""
        strings = []
        for i, value in enumerate(obj_list):
            new_path = f"{path}[{i}]" if path else f"[{i}]"
            strings.extend(self._extract_yaml_strings(value, new_path))
        return strings

    def _extract_from_odt_sync(self, content: bytes) -> str:
        """Синхронное извлечение текста из ODT."""
        if not load:
            raise ImportError("odfpy не установлен")

        temp_file_path = None
        try:
            with tempfile.NamedTemporaryFile(suffix=".odt", delete=False) as temp_file:
                temp_file.write(content)
                temp_file_path = temp_file.name

            doc = load(temp_file_path)
            text_parts = []

            # Извлечение всех текстовых элементов
            for p in doc.getElementsByType(P):
                text = extractText(p)
                if text.strip():
                    text_parts.append(text)

            return "\n".join(text_parts)

        except Exception as e:
            logger.error(f"Ошибка при обработке ODT: {str(e)}")
            raise ValueError(f"Error processing ODT: {str(e)}")
        finally:
            # Гарантированное удаление временного файла
            if temp_file_path and os.path.exists(temp_file_path):
                try:
                    os.unlink(temp_file_path)
                except OSError as e:
                    logger.warning(
                        f"Не удалось удалить временный файл {temp_file_path}: {str(e)}"
                    )

    def _extract_from_epub_sync(self, content: bytes) -> str:
        """Синхронное извлечение текста из EPUB."""
        if not BeautifulSoup:
            raise ImportError("beautifulsoup4 не установлен")

        try:
            text_parts = []
            extracted_size = 0

            with zipfile.ZipFile(io.BytesIO(content), "r") as zip_ref:
                for file_info in zip_ref.infolist():
                    if self._should_stop_epub_extraction(
                        extracted_size, file_info.file_size
                    ):
                        break

                    if self._is_epub_html_file(file_info.filename):
                        text, new_size = self._extract_epub_html_text(
                            zip_ref, file_info
                        )
                        if text:
                            text_parts.append(text)
                        extracted_size += new_size

            return "\n\n".join(text_parts)

        except Exception as e:
            logger.error(f"Ошибка при обработке EPUB: {str(e)}")
            raise ValueError(f"Error processing EPUB: {str(e)}")

    def _should_stop_epub_extraction(self, current_size: int, file_size: int) -> bool:
        """Проверка лимита размера для EPUB."""
        if current_size + file_size > settings.MAX_EXTRACTED_SIZE:
            logger.warning("Достигнут лимит размера распакованного содержимого EPUB")
            return True
        return False

    def _is_epub_html_file(self, filename: str) -> bool:
        """Проверка является ли файл HTML для EPUB."""
        return filename.endswith((".html", ".xhtml", ".htm"))

    def _extract_epub_html_text(self, zip_ref, file_info) -> tuple:
        """Извлечение текста из HTML файла EPUB."""
        try:
            html_content = zip_ref.read(file_info.filename)
            html_text = html_content.decode("utf-8", errors="replace")

            # Парсинг HTML
            soup = BeautifulSoup(html_text, "html.parser")

            # Удаление script и style тегов
            for script in soup(["script", "style"]):
                script.decompose()

            # Извлечение текста
            text = soup.get_text()
            return text.strip() if text.strip() else None, file_info.file_size

        except Exception as e:
            logger.warning(f"Ошибка при обработке файла {file_info.filename}: {e}")
            return None, 0

    def _extract_from_eml_sync(self, content: bytes) -> str:
        """Синхронное извлечение текста из EML."""
        import email

        try:
            msg_text = self._decode_eml_content(content)
            msg = email.message_from_string(msg_text)
            text_parts = []

            # Извлечение заголовков
            text_parts.extend(self._extract_eml_headers(msg))
            text_parts.append("---")

            # Извлечение тела письма
            if msg.is_multipart():
                text_parts.extend(self._extract_eml_body_multipart(msg))
            else:
                text_parts.extend(self._extract_eml_body_simple(msg))

            return (
                "\n".join(text_parts)
                if text_parts
                else "Не удалось извлечь читаемый текст из EML файла"
            )

        except Exception as e:
            logger.error(f"Ошибка при обработке EML: {str(e)}")
            raise ValueError(f"Error processing EML: {str(e)}")

    def _decode_eml_content(self, content: bytes) -> str:
        """Декодирование содержимого EML файла."""
        for encoding in ["utf-8", "cp1251", "latin-1"]:
            try:
                return content.decode(encoding)
            except UnicodeDecodeError:
                continue
        return content.decode("utf-8", errors="replace")

    def _extract_eml_headers(self, msg) -> list:
        """Извлечение заголовков из EML."""
        from email.header import decode_header

        text_parts = []
        headers = ["From", "To", "Subject", "Date"]

        for header in headers:
            value = msg.get(header)
            if value:
                decoded_value = self._decode_eml_header(value, decode_header)
                text_parts.append(f"{header}: {decoded_value}")

        return text_parts

    def _decode_eml_header(self, value: str, decode_header_func) -> str:
        """Декодирование заголовка EML."""
        decoded_parts = decode_header_func(value)
        decoded_value = ""

        for part, encoding in decoded_parts:
            if isinstance(part, bytes):
                if encoding:
                    decoded_value += part.decode(encoding)
                else:
                    decoded_value += part.decode("utf-8", errors="replace")
            else:
                decoded_value += part

        return decoded_value

    def _extract_eml_body_multipart(self, msg) -> list:
        """Извлечение тела многочастного EML письма."""
        text_parts = []

        for part in msg.walk():
            content_type = part.get_content_type()
            if content_type in ["text/plain", "text/html"]:
                try:
                    body_text = self._extract_eml_part_text(part, content_type)
                    if body_text and body_text.strip():
                        text_parts.append(body_text)
                except Exception as e:
                    logger.warning(f"Ошибка при обработке части письма: {e}")

        return text_parts

    def _extract_eml_body_simple(self, msg) -> list:
        """Извлечение тела простого EML письма."""
        text_parts = []

        try:
            payload = msg.get_payload(decode=True)
            if payload:
                charset = msg.get_content_charset() or "utf-8"
                body_text = self._decode_payload(payload, charset)
                if body_text.strip():
                    text_parts.append(body_text)
        except Exception as e:
            logger.warning(f"Ошибка при обработке тела письма: {e}")

        return text_parts

    def _extract_eml_part_text(self, part, content_type: str) -> str:
        """Извлечение текста из части EML."""
        payload = part.get_payload(decode=True)
        if not payload:
            return ""

        charset = part.get_content_charset() or "utf-8"
        body_text = self._decode_payload(payload, charset)

        # Обработка HTML
        if content_type == "text/html" and BeautifulSoup:
            soup = BeautifulSoup(body_text, "html.parser")
            body_text = soup.get_text()

        return body_text

    def _decode_payload(self, payload: bytes, charset: str) -> str:
        """Декодирование payload с обработкой ошибок."""
        try:
            return payload.decode(charset)
        except UnicodeDecodeError:
            return payload.decode("utf-8", errors="replace")

    def _extract_from_msg_sync(self, content: bytes) -> str:
        """Синхронное извлечение текста из MSG."""
        try:
            text_parts = []

            # Извлечение UTF-16 текста
            text_parts.extend(self._extract_utf16_text_from_msg(content))

            # Альтернативный подход - поиск ASCII текста
            text_parts.extend(self._extract_ascii_text_from_msg(content, text_parts))

            return (
                "\n".join(text_parts)
                if text_parts
                else "Не удалось извлечь читаемый текст из MSG файла"
            )

        except Exception as e:
            logger.error(f"Ошибка при обработке MSG: {str(e)}")
            raise ValueError(f"Error processing MSG: {str(e)}")

    def _extract_utf16_text_from_msg(self, content: bytes) -> list:
        """Извлечение UTF-16 текста из MSG файла."""
        text_parts = []
        try:
            text = content.decode("utf-16le", errors="ignore")
            lines = text.split("\n")
            clean_lines = self._clean_msg_lines(lines)
            unique_lines = self._filter_unique_lines(clean_lines, min_length=5)
            text_parts.extend(unique_lines)
        except Exception as e:
            logger.warning(f"Ошибка при декодировании UTF-16: {e}")
        return text_parts

    def _clean_msg_lines(self, lines: list) -> list:
        """Очистка строк MSG от управляющих символов."""
        clean_lines = []
        for line in lines:
            # Убираем нулевые байты и управляющие символы
            clean_line = "".join(
                char for char in line if ord(char) >= 32 or char in "\t\n\r"
            )
            clean_line = clean_line.strip()

            # Пропускаем слишком короткие или бессмысленные строки
            if self._is_valid_msg_line(clean_line):
                clean_lines.append(clean_line)

        return clean_lines

    def _is_valid_msg_line(self, line: str) -> bool:
        """Проверка валидности строки MSG."""
        return (
            len(line) > 3
            and not line.startswith(("_", "\x00"))
            and any(c.isalpha() for c in line)
        )

    def _filter_unique_lines(self, lines: list, min_length: int = 5) -> list:
        """Фильтрация уникальных строк с минимальной длиной."""
        unique_lines = []
        seen = set()
        for line in lines:
            if line not in seen and len(line) > min_length:
                unique_lines.append(line)
                seen.add(line)
        return unique_lines

    def _extract_ascii_text_from_msg(
        self, content: bytes, existing_text_parts: list
    ) -> list:
        """Извлечение ASCII текста из MSG файла."""
        text_parts = []
        try:
            ascii_text = content.decode("ascii", errors="ignore")
            lines = ascii_text.split("\n")

            for line in lines:
                clean_line = line.strip()
                if self._is_valid_ascii_line(clean_line, existing_text_parts):
                    text_parts.append(clean_line)
        except Exception as e:
            logger.warning(f"Ошибка при извлечении ASCII: {e}")
        return text_parts

    def _is_valid_ascii_line(self, line: str, existing_parts: list) -> bool:
        """Проверка валидности ASCII строки."""
        return (
            len(line) > 10
            and any(c.isalpha() for c in line)
            and line not in existing_parts
        )

    def _safe_tesseract_ocr(self, image, temp_image_path: str = None) -> str:
        """
        Безопасный вызов Tesseract с ограничениями ресурсов.

        Args:
            image: PIL Image объект
            temp_image_path: Путь к временному файлу (если None, создается автоматически)

        Returns:
            str: Распознанный текст
        """
        import os
        import tempfile

        from .config import settings
        from .utils import run_subprocess_with_limits

        temp_file_created = False

        try:
            # Если путь к временному файлу не указан, создаем его
            if temp_image_path is None:
                with tempfile.NamedTemporaryFile(
                    suffix=".png", delete=False
                ) as temp_file:
                    temp_image_path = temp_file.name
                    temp_file_created = True

                # Сохраняем изображение во временный файл
                # Конвертируем в RGB для совместимости с PNG
                if image.mode in ("RGBA", "LA", "P"):
                    image = image.convert("RGB")
                image.save(temp_image_path, "PNG")

            # Создаем временный файл для вывода
            with tempfile.NamedTemporaryFile(
                suffix=".txt", delete=False
            ) as output_file:
                output_path = output_file.name

            try:
                # Вызываем tesseract через безопасную функцию
                result = run_subprocess_with_limits(
                    command=[
                        "tesseract",
                        temp_image_path,
                        output_path.replace(".txt", ""),
                        "-l",
                        self.ocr_languages,
                    ],
                    timeout=30,
                    memory_limit=settings.MAX_TESSERACT_MEMORY,
                    capture_output=True,
                    text=True,
                )

                if result.returncode != 0:
                    logger.warning(
                        f"Tesseract завершился с кодом {result.returncode}: {result.stderr}"
                    )
                    return ""

                # Читаем результат OCR
                if os.path.exists(output_path):
                    with open(output_path, "r", encoding="utf-8") as f:
                        return f.read().strip()
                else:
                    logger.warning("Файл результата OCR не найден")
                    return ""

            finally:
                # Удаляем временный файл вывода
                try:
                    if os.path.exists(output_path):
                        os.unlink(output_path)
                except Exception as e:
                    logger.warning(
                        f"Не удалось удалить временный файл {output_path}: {e}"
                    )

        except subprocess.TimeoutExpired:
            logger.error("Tesseract OCR timeout")
            return ""
        except MemoryError as e:
            logger.error(f"Tesseract превысил лимит памяти: {str(e)}")
            return ""
        except Exception as e:
            logger.error(f"Ошибка при OCR: {str(e)}")
            return ""
        finally:
            # Удаляем временный файл изображения, если мы его создали
            if (
                temp_file_created
                and temp_image_path
                and os.path.exists(temp_image_path)
            ):
                try:
                    os.unlink(temp_image_path)
                except Exception as e:
                    logger.warning(
                        f"Не удалось удалить временный файл {temp_image_path}: {e}"
                    )

    def _extract_from_image_sync(self, content: bytes) -> str:
        """Синхронный OCR изображения."""
        if not Image:
            raise ImportError("PIL не установлен")

        try:
            # Валидация изображения для предотвращения DoS атак
            from .utils import validate_image_for_ocr  # noqa: F811

            is_valid, error_message = validate_image_for_ocr(content)
            if not is_valid:
                logger.warning(f"Изображение не прошло валидацию: {error_message}")
                raise ValueError(f"Image validation failed: {error_message}")

            image = Image.open(io.BytesIO(content))

            # Безопасный OCR с ограничениями ресурсов
            text = self._safe_tesseract_ocr(image)
            return text

        except Exception as e:
            logger.error(f"Ошибка при OCR изображения: {str(e)}")
            raise ValueError(f"Error processing image: {str(e)}")

    def _check_mime_type(self, content: bytes, filename: str) -> bool:
        """Проверка MIME-типа файла для предотвращения подделки расширений."""
        import mimetypes

        try:
            # Определяем MIME-тип по содержимому (первые байты)
            mime_signatures = {
                b"\x50\x4b\x03\x04": [
                    "application/zip",
                    "application/epub+zip",
                    "application/vnd.openxmlformats",
                ],
                b"\x50\x4b\x07\x08": ["application/zip", "application/epub+zip"],
                b"\x50\x4b\x05\x06": ["application/zip", "application/epub+zip"],
                b"%PDF": ["application/pdf"],
                b"\xd0\xcf\x11\xe0": [
                    "application/msword",
                    "application/vnd.ms-excel",
                    "application/vnd.ms-powerpoint",
                ],
                b"\x89PNG": ["image/png"],
                b"\xff\xd8\xff": ["image/jpeg"],
                b"GIF8": ["image/gif"],
                b"BM": ["image/bmp"],
                b"II*\x00": ["image/tiff"],
                b"MM\x00*": ["image/tiff"],
                b"<!DOCTYPE": ["text/html"],
                b"<html": ["text/html"],
                b"<?xml": ["text/xml", "application/xml"],
            }

            # Проверяем сигнатуру файла
            file_start = content[:10]
            detected_mime = None

            for signature, mime_types in mime_signatures.items():
                if file_start.startswith(signature):
                    detected_mime = mime_types[0]
                    break

            # Определяем ожидаемый MIME-тип по расширению
            expected_mime, _ = mimetypes.guess_type(filename)

            # Если не можем определить MIME-тип, разрешаем
            if not detected_mime or not expected_mime:
                return True

            # Проверяем соответствие
            return detected_mime in mime_signatures.get(file_start[:4], [expected_mime])

        except Exception as e:
            logger.warning(f"Ошибка при проверке MIME-типа: {str(e)}")
            return True  # В случае ошибки разрешаем обработку

    def _extract_from_archive(
        self, content: bytes, filename: str, nesting_level: int = 0
    ) -> List[Dict[str, Any]]:
        """Безопасное извлечение файлов из архива."""
        # Проверка глубины вложенности
        if nesting_level >= settings.MAX_ARCHIVE_NESTING:
            logger.warning(
                f"Превышена максимальная глубина вложенности архивов: {filename}"
            )
            raise ValueError("Maximum archive nesting level exceeded")

        # Проверка размера архива
        if len(content) > settings.MAX_ARCHIVE_SIZE:
            logger.warning(f"Архив {filename} слишком большой: {len(content)} байт")
            raise ValueError("Archive size exceeds maximum allowed size")

        extension = get_file_extension(filename)
        logger.info(
            f"Обработка архива {filename} (тип: {extension}, размер: {len(content)} байт)"
        )

        extracted_files = []

        # Создаем временную директорию для безопасной работы
        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)
            archive_path = temp_path / f"archive_{int(time.time())}.{extension}"

            try:
                # Записываем архив во временный файл
                with open(archive_path, "wb") as f:
                    f.write(content)

                # Извлекаем файлы в зависимости от типа архива
                extract_dir = temp_path / "extracted"
                extract_dir.mkdir(exist_ok=True)

                if extension == "zip":
                    extracted_files = self._extract_zip_files(
                        archive_path, extract_dir, filename, nesting_level
                    )
                elif extension in [
                    "tar",
                    "gz",
                    "bz2",
                    "xz",
                    "tar.gz",
                    "tar.bz2",
                    "tar.xz",
                    "tgz",
                    "tbz2",
                    "txz",
                ]:
                    extracted_files = self._extract_tar_files(
                        archive_path, extract_dir, filename, nesting_level
                    )
                elif extension == "rar":
                    extracted_files = self._extract_rar_files(
                        archive_path, extract_dir, filename, nesting_level
                    )
                elif extension == "7z":
                    extracted_files = self._extract_7z_files(
                        archive_path, extract_dir, filename, nesting_level
                    )
                else:
                    raise ValueError(f"Unsupported archive format: {extension}")

                logger.info(
                    f"Успешно обработано {len(extracted_files)} файлов из архива {filename}"
                )
                return extracted_files

            except Exception as e:
                logger.error(f"Ошибка при обработке архива {filename}: {str(e)}")
                raise ValueError(f"Error processing archive: {str(e)}")

    def _extract_zip_files(
        self,
        archive_path: Path,
        extract_dir: Path,
        archive_name: str,
        nesting_level: int,
    ) -> List[Dict[str, Any]]:
        """Извлечение файлов из ZIP-архива."""
        try:
            with zipfile.ZipFile(archive_path, "r") as zip_ref:
                self._validate_zip_size(zip_ref)
                return self._process_zip_files(
                    zip_ref, extract_dir, archive_name, nesting_level
                )
        except zipfile.BadZipFile:
            raise ValueError("Invalid ZIP file")

    def _validate_zip_size(self, zip_ref) -> None:
        """Проверка размера файлов в ZIP-архиве для защиты от zip bomb."""
        total_size = 0
        for info in zip_ref.infolist():
            if not info.is_dir():
                total_size += info.file_size
                if total_size > settings.MAX_EXTRACTED_SIZE:
                    raise ValueError(
                        "Extracted files size exceeds maximum allowed size (zip bomb protection)"
                    )

    def _process_zip_files(
        self, zip_ref, extract_dir: Path, archive_name: str, nesting_level: int
    ) -> List[Dict[str, Any]]:
        """Обработка всех файлов в ZIP-архиве."""
        extracted_files = []

        for info in zip_ref.infolist():
            if info.is_dir():
                continue

            file_result = self._extract_single_zip_file(
                info, zip_ref, extract_dir, archive_name, nesting_level
            )
            if file_result:
                extracted_files.extend(file_result)

        return extracted_files

    def _extract_single_zip_file(
        self, info, zip_ref, extract_dir: Path, archive_name: str, nesting_level: int
    ) -> List[Dict[str, Any]]:
        """Извлечение и обработка одного файла из ZIP-архива."""
        # Санитизируем имя файла
        safe_filename = self._sanitize_archive_filename(info.filename)
        if not safe_filename:
            return []

        # Фильтруем системные файлы
        if self._is_system_file(safe_filename):
            return []

        # Создаем безопасный путь для извлечения
        safe_path = extract_dir / safe_filename
        safe_path.parent.mkdir(parents=True, exist_ok=True)

        try:
            # Извлекаем файл
            with zip_ref.open(info) as source, open(safe_path, "wb") as target:
                shutil.copyfileobj(source, target)

            # Обрабатываем файл
            file_content = safe_path.read_bytes()
            return (
                self._process_extracted_file(
                    file_content,
                    safe_filename,
                    safe_path.name,
                    archive_name,
                    nesting_level,
                )
                or []
            )

        except Exception as e:
            logger.warning(
                f"Ошибка при обработке файла {safe_filename} из архива {archive_name}: {str(e)}"
            )
            return []

    def _extract_tar_files(
        self,
        archive_path: Path,
        extract_dir: Path,
        archive_name: str,
        nesting_level: int,
    ) -> List[Dict[str, Any]]:
        """Извлечение файлов из TAR-архива."""
        extracted_files = []
        total_size = 0

        try:
            with tarfile.open(archive_path, "r:*") as tar_ref:
                # Проверяем размер распакованных файлов
                for member in tar_ref.getmembers():
                    if member.isfile():
                        total_size += member.size

                        if total_size > settings.MAX_EXTRACTED_SIZE:
                            raise ValueError(
                                "Extracted files size exceeds maximum allowed size (tar bomb protection)"
                            )

                # Извлекаем файлы
                for member in tar_ref.getmembers():
                    if not member.isfile():
                        continue

                    # Санитизируем имя файла
                    safe_filename = self._sanitize_archive_filename(member.name)
                    if not safe_filename:
                        continue

                    # Фильтруем системные файлы
                    if self._is_system_file(safe_filename):
                        continue

                    # Создаем безопасный путь для извлечения
                    safe_path = extract_dir / safe_filename
                    safe_path.parent.mkdir(parents=True, exist_ok=True)

                    try:
                        # Извлекаем файл
                        with (
                            tar_ref.extractfile(member) as source,
                            open(safe_path, "wb") as target,
                        ):
                            if source:
                                shutil.copyfileobj(source, target)

                        # Обрабатываем файл
                        file_content = safe_path.read_bytes()
                        file_result = self._process_extracted_file(
                            file_content,
                            safe_filename,
                            safe_path.name,
                            archive_name,
                            nesting_level,
                        )

                        if file_result:
                            extracted_files.extend(file_result)

                    except Exception as e:
                        logger.warning(
                            f"Ошибка при обработке файла {safe_filename} из архива {archive_name}: {str(e)}"
                        )
                        continue

        except tarfile.TarError:
            raise ValueError("Invalid TAR file")

        return extracted_files

    def _extract_rar_files(
        self,
        archive_path: Path,
        extract_dir: Path,
        archive_name: str,
        nesting_level: int,
    ) -> List[Dict[str, Any]]:
        """Извлечение файлов из RAR-архива."""
        if not rarfile:
            raise ValueError("RAR support not available. Install rarfile library.")

        extracted_files = []
        total_size = 0

        try:
            with rarfile.RarFile(archive_path, "r") as rar_ref:
                # Проверяем размер распакованных файлов
                for info in rar_ref.infolist():
                    if info.is_dir():
                        continue
                    total_size += info.file_size

                    if total_size > settings.MAX_EXTRACTED_SIZE:
                        raise ValueError(
                            "Extracted files size exceeds maximum allowed size (rar bomb protection)"
                        )

                # Извлекаем файлы
                for info in rar_ref.infolist():
                    if info.is_dir():
                        continue

                    # Санитизируем имя файла
                    safe_filename = self._sanitize_archive_filename(info.filename)
                    if not safe_filename:
                        continue

                    # Фильтруем системные файлы
                    if self._is_system_file(safe_filename):
                        continue

                    # Создаем безопасный путь для извлечения
                    safe_path = extract_dir / safe_filename
                    safe_path.parent.mkdir(parents=True, exist_ok=True)

                    try:
                        # Извлекаем файл
                        with (
                            rar_ref.open(info) as source,
                            open(safe_path, "wb") as target,
                        ):
                            shutil.copyfileobj(source, target)

                        # Обрабатываем файл
                        file_content = safe_path.read_bytes()
                        file_result = self._process_extracted_file(
                            file_content,
                            safe_filename,
                            safe_path.name,
                            archive_name,
                            nesting_level,
                        )

                        if file_result:
                            extracted_files.extend(file_result)

                    except Exception as e:
                        logger.warning(
                            f"Ошибка при обработке файла {safe_filename} из архива {archive_name}: {str(e)}"
                        )
                        continue

        except rarfile.RarError:
            raise ValueError("Invalid RAR file")

        return extracted_files

    def _extract_7z_files(
        self,
        archive_path: Path,
        extract_dir: Path,
        archive_name: str,
        nesting_level: int,
    ) -> List[Dict[str, Any]]:
        """Извлечение файлов из 7Z-архива."""
        if not py7zr:
            raise ValueError("7Z support not available. Install py7zr library.")

        extracted_files = []
        total_size = 0

        try:
            with py7zr.SevenZipFile(archive_path, "r") as sz_ref:
                # Проверяем размер распакованных файлов
                for info in sz_ref.list():
                    if info.is_dir:
                        continue
                    total_size += info.uncompressed

                    if total_size > settings.MAX_EXTRACTED_SIZE:
                        raise ValueError(
                            "Extracted files size exceeds maximum allowed size (7z bomb protection)"
                        )

                # Извлекаем файлы
                sz_ref.extractall(extract_dir)

                # Обрабатываем извлеченные файлы
                for root, _dirs, files in os.walk(extract_dir):
                    for file in files:
                        file_path = Path(root) / file
                        relative_path = file_path.relative_to(extract_dir)

                        # Санитизируем имя файла
                        safe_filename = self._sanitize_archive_filename(
                            str(relative_path)
                        )
                        if not safe_filename:
                            continue

                        # Фильтруем системные файлы
                        if self._is_system_file(safe_filename):
                            continue

                        try:
                            # Обрабатываем файл
                            file_content = file_path.read_bytes()
                            file_result = self._process_extracted_file(
                                file_content,
                                safe_filename,
                                file_path.name,
                                archive_name,
                                nesting_level,
                            )

                            if file_result:
                                extracted_files.extend(file_result)

                        except Exception as e:
                            logger.warning(
                                f"Ошибка при обработке файла {safe_filename} из архива {archive_name}: {str(e)}"
                            )
                            continue

        except py7zr.Bad7zFile:
            raise ValueError("Invalid 7Z file")

        return extracted_files

    def _process_extracted_file(
        self,
        content: bytes,
        filename: str,
        basename: str,
        archive_name: str,
        nesting_level: int,
    ) -> Optional[List[Dict[str, Any]]]:
        """Обработка извлеченного файла."""
        try:
            # Если файл является архивом, рекурсивно обрабатываем его
            if is_archive_format(basename, settings.SUPPORTED_FORMATS):
                return self._extract_from_archive(content, basename, nesting_level + 1)

            # Если файл поддерживается, извлекаем текст
            if is_supported_format(basename, settings.SUPPORTED_FORMATS):
                extension = get_file_extension(basename)
                text = self._extract_text_by_format(content, extension, basename)

                return [
                    {
                        "filename": basename,
                        "path": f"{archive_name}/{filename}",
                        "size": len(content),
                        "type": extension,
                        "text": text.strip() if text else "",
                    }
                ]

            return None

        except Exception as e:
            logger.warning(f"Ошибка при обработке файла {filename}: {str(e)}")
            return None

    def _sanitize_archive_filename(self, filename: str) -> str:
        """Санитизация имени файла из архива."""
        if not filename:
            return ""

        # Удаляем опасные пути
        filename = filename.replace("..", "").replace("\\", "/").strip("/")

        # Проверяем на абсолютные пути
        if filename.startswith("/"):
            filename = filename[1:]

        # Удаляем пустые части пути
        parts = [part for part in filename.split("/") if part and part != "."]

        if not parts:
            return ""

        return "/".join(parts)

    def _is_system_file(self, filename: str) -> bool:
        """Проверка, является ли файл системным."""
        system_files = [
            ".DS_Store",
            "Thumbs.db",
            ".git/",
            ".svn/",
            ".hg/",
            "__MACOSX/",
            ".localized",
            "desktop.ini",
            "folder.ini",
        ]

        filename_lower = filename.lower()
        for system_file in system_files:
            if system_file in filename_lower:
                return True

        return False

    def _ocr_from_pdf_image_sync(self, page, img_info) -> str:
        """Синхронный OCR изображения из PDF."""
        if not Image:
            return ""

        try:
            # Получаем координаты изображения
            x0, y0, x1, y1 = (
                img_info["x0"],
                img_info["y0"],
                img_info["x1"],
                img_info["y1"],
            )

            # Проверяем разумность размеров области
            width = abs(x1 - x0)
            height = abs(y1 - y0)

            # Ограничиваем размер области для предотвращения DoS
            max_dimension = 5000  # максимальный размер по любой оси
            if width > max_dimension or height > max_dimension:
                logger.warning(f"Область изображения слишком большая: {width}x{height}")
                return ""

            # Обрезаем область изображения из всей страницы
            cropped_bbox = (x0, y0, x1, y1)
            cropped_page = page.crop(cropped_bbox)

            # Конвертируем обрезанную область в изображение с высоким разрешением
            img_pil = cropped_page.to_image(resolution=300)

            # Безопасный OCR с ограничениями ресурсов
            text = self._safe_tesseract_ocr(img_pil.original)

            return text

        except Exception as e:
            logger.warning(f"Ошибка OCR изображения: {str(e)}")
            # Альтернативный подход - рендерим всю страницу и обрезаем
            try:
                # Конвертируем всю страницу в изображение PIL
                page_image = page.to_image(resolution=300)
                pil_image = page_image.original  # Получаем PIL изображение

                # Вычисляем координаты в пикселях (учитывая resolution=300)
                scale = 300 / 72  # PDF обычно 72 DPI, мы рендерим в 300 DPI
                pixel_bbox = (
                    int(x0 * scale),
                    int(y0 * scale),
                    int(x1 * scale),
                    int(y1 * scale),
                )

                # Проверяем разумность размеров в пикселях
                pixel_width = abs(pixel_bbox[2] - pixel_bbox[0])
                pixel_height = abs(pixel_bbox[3] - pixel_bbox[1])

                if pixel_width * pixel_height > 25000000:  # 25MP максимум
                    logger.warning(
                        f"Область изображения слишком большая: {pixel_width}x{pixel_height} пикселей"
                    )
                    return ""

                # Обрезаем область изображения
                cropped_img = pil_image.crop(pixel_bbox)

                # Безопасный OCR с ограничениями ресурсов
                text = self._safe_tesseract_ocr(cropped_img)

                return text

            except Exception as e2:
                logger.warning(
                    f"Альтернативная попытка OCR также не удалась: {str(e2)}"
                )
                return ""

    # Веб-экстракция (новое в v1.10.0)

    def _extract_page_with_playwright(
        self,
        url: str,
        user_agent: Optional[str] = None,
        extraction_options: Optional[Any] = None,
    ) -> tuple[str, str]:
        """
        Извлечение HTML контента страницы с помощью Playwright (с поддержкой JS, обновлено в v1.10.2).

        Args:
            url: URL страницы
            user_agent: Пользовательский User-Agent
            extraction_options: Настройки извлечения

        Returns:
            tuple[str, str]: (html_content, final_url)
        """
        if not sync_playwright:
            raise ValueError("Playwright не установлен")

        # Определяем настройки с учетом переданных параметров или значений по умолчанию
        web_page_timeout = (
            extraction_options.web_page_timeout
            if extraction_options and extraction_options.web_page_timeout is not None
            else settings.WEB_PAGE_TIMEOUT
        )

        js_render_timeout = (
            extraction_options.js_render_timeout
            if extraction_options and extraction_options.js_render_timeout is not None
            else settings.JS_RENDER_TIMEOUT
        )

        web_page_delay = (
            extraction_options.web_page_delay
            if extraction_options and extraction_options.web_page_delay is not None
            else settings.WEB_PAGE_DELAY
        )

        enable_lazy_loading_wait = (
            extraction_options.enable_lazy_loading_wait
            if extraction_options
            and extraction_options.enable_lazy_loading_wait is not None
            else settings.ENABLE_LAZY_LOADING_WAIT
        )

        html_content = ""
        final_url = url

        with sync_playwright() as p:
            # Запускаем Chromium (установлен в Dockerfile)
            browser = p.chromium.launch(
                headless=True,
                args=[
                    "--no-sandbox",
                    "--disable-dev-shm-usage",
                    "--disable-gpu",
                    "--disable-extensions",
                    "--disable-web-security",  # для обхода CORS при локальной разработке
                ],
            )

            try:
                # Определяем, включать ли JavaScript
                enable_javascript = (
                    extraction_options.enable_javascript
                    if extraction_options
                    and extraction_options.enable_javascript is not None
                    else settings.ENABLE_JAVASCRIPT
                )

                context = browser.new_context(
                    user_agent=user_agent or settings.DEFAULT_USER_AGENT,
                    viewport={"width": 1280, "height": 720},
                    java_script_enabled=enable_javascript,  # Правильная настройка отключения JS
                )

                page = context.new_page()

                # Устанавливаем таймауты для защиты от DoS-атак
                page.set_default_timeout(web_page_timeout * 1000)  # в миллисекундах

                # Дополнительная защита от DoS: ограничиваем время выполнения JavaScript
                page.set_default_navigation_timeout(web_page_timeout * 1000)

                # Переходим на страницу
                logger.info(
                    f"Загрузка страницы с Playwright: {url} (JS: {'включен' if enable_javascript else 'отключен'})"
                )
                response = page.goto(url, wait_until="domcontentloaded")

                if not response.ok:
                    raise ValueError(f"HTTP {response.status}: {response.status_text}")

                final_url = page.url

                # Ждем дополнительной загрузки JS (если включено)
                if enable_javascript:
                    logger.info(f"Ожидание JS-рендеринга ({js_render_timeout}s)...")

                    # Ждем загрузки сети с защитой от DoS атак
                    try:
                        # Ограничиваем время выполнения для защиты от ресурсоемких скриптов
                        page.wait_for_load_state(
                            "networkidle",
                            timeout=min(
                                js_render_timeout * 1000, 15000
                            ),  # не более 15 сек
                        )
                    except Exception as e:
                        logger.warning(
                            f"Таймаут ожидания сети (защита от DoS): {str(e)}"
                        )

                    # Обработка lazy loading с защитой от бесконечности
                    if enable_lazy_loading_wait:
                        self._safe_scroll_for_lazy_loading(page, extraction_options)

                    # Дополнительная задержка для завершения JS
                    import time

                    time.sleep(web_page_delay)

                # Получаем финальный HTML
                html_content = page.content()
                logger.info(f"HTML получен, размер: {len(html_content)} символов")

            finally:
                browser.close()

        return html_content, final_url

    def _safe_scroll_for_lazy_loading(
        self, page, extraction_options: Optional[Any] = None
    ) -> None:
        """
        Безопасный скролл страницы для активации lazy loading с защитой от бесконечности (обновлено в v1.10.2).

        Args:
            page: Playwright page объект
            extraction_options: Настройки извлечения
        """
        try:
            logger.info("Начинаем безопасный скролл для активации lazy loading...")

            # Определяем максимальное количество попыток скролла
            max_scroll_attempts = (
                extraction_options.max_scroll_attempts
                if extraction_options
                and extraction_options.max_scroll_attempts is not None
                else settings.MAX_SCROLL_ATTEMPTS
            )

            # Получаем начальную высоту страницы
            initial_height = page.evaluate("document.body.scrollHeight")
            logger.info(f"Начальная высота страницы: {initial_height}px")

            scroll_attempts = 0
            last_height = initial_height
            stable_count = 0  # Счетчик стабильных измерений

            while scroll_attempts < max_scroll_attempts:
                scroll_attempts += 1
                logger.info(f"Попытка скролла {scroll_attempts}/{max_scroll_attempts}")

                # Плавный скролл до конца страницы
                page.evaluate("window.scrollTo(0, document.body.scrollHeight)")

                # Ждем небольшую задержку для загрузки контента
                import time

                time.sleep(1)

                # Проверяем новую высоту
                new_height = page.evaluate("document.body.scrollHeight")
                logger.info(f"Новая высота страницы: {new_height}px")

                # Если высота не изменилась, увеличиваем счетчик стабильности
                if new_height == last_height:
                    stable_count += 1
                    logger.info(f"Высота стабильна, счетчик: {stable_count}")

                    # Если высота стабильна уже 2 раза подряд - прекращаем
                    if stable_count >= 2:
                        logger.info(
                            "Высота страницы стабилизировалась, завершаем скролл"
                        )
                        break
                else:
                    # Высота изменилась, сбрасываем счетчик
                    stable_count = 0
                    last_height = new_height

                # Дополнительная проверка: если страница выросла слишком сильно, прекращаем
                if new_height > initial_height * 10:  # Если страница выросла в 10 раз
                    logger.warning(
                        "Страница выросла подозрительно сильно, возможно бесконечный скролл"
                    )
                    break

            # Возвращаемся в начало страницы
            page.evaluate("window.scrollTo(0, 0)")
            logger.info("Скролл завершен, возвращены в начало страницы")

        except Exception as e:
            logger.warning(f"Ошибка при скролле для lazy loading: {str(e)}")

    def _determine_content_type(
        self,
        url: str,
        user_agent: Optional[str] = None,
        extraction_options: Optional[Any] = None,
    ) -> tuple[str, str]:
        """
        Определение типа контента через HEAD запрос (новое в v1.10.3).

        Args:
            url: URL для проверки
            user_agent: Пользовательский User-Agent
            extraction_options: Настройки извлечения

        Returns:
            tuple[str, str]: (content_type, final_url) - тип контента и финальный URL после редиректов
        """
        if not requests:
            raise ValueError(
                "requests library not available for content type determination"
            )

        # Определяем настройки с учетом переданных параметров или значений по умолчанию
        head_timeout = (
            extraction_options.web_page_timeout
            if extraction_options and extraction_options.web_page_timeout is not None
            else settings.HEAD_REQUEST_TIMEOUT
        )

        follow_redirects = (
            extraction_options.follow_redirects
            if extraction_options and extraction_options.follow_redirects is not None
            else True
        )

        max_redirects = (
            extraction_options.max_redirects
            if extraction_options and extraction_options.max_redirects is not None
            else 5
        )

        # Установка User-Agent и заголовков
        headers = {
            "User-Agent": user_agent or settings.DEFAULT_USER_AGENT,
            "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8",
            "Accept-Language": "ru,en;q=0.5",
            "Accept-Encoding": "gzip, deflate",
            "Connection": "keep-alive",
        }

        # Создание сессии для следования редиректам
        session = requests.Session()
        session.headers.update(headers)

        try:
            logger.info(f"Выполняю HEAD запрос для определения типа контента: {url}")

            response = session.head(
                url, timeout=head_timeout, allow_redirects=follow_redirects, stream=True
            )

            if follow_redirects and len(response.history) > max_redirects:
                logger.warning(
                    f"Превышено максимальное количество редиректов ({max_redirects})"
                )

            response.raise_for_status()

            content_type = response.headers.get("content-type", "").lower()
            final_url = response.url

            logger.info(f"Определен Content-Type: {content_type} для URL: {final_url}")

            return content_type, final_url

        except Exception as e:
            logger.warning(f"Ошибка HEAD запроса: {str(e)}, попробуем GET запрос")
            # Fallback: делаем GET запрос но читаем только заголовки
            try:
                response = session.get(
                    url,
                    timeout=head_timeout,
                    allow_redirects=follow_redirects,
                    stream=True,
                )
                response.raise_for_status()

                content_type = response.headers.get("content-type", "").lower()
                final_url = response.url

                # Закрываем соединение, не читая тело
                response.close()

                logger.info(
                    f"Определен Content-Type через GET: {content_type} для URL: {final_url}"
                )
                return content_type, final_url

            except Exception as get_error:
                logger.error(f"Ошибка при определении типа контента: {str(get_error)}")
                raise ValueError(f"Unable to determine content type: {str(get_error)}")
        finally:
            session.close()

    def _is_html_content(self, content_type: str, url: str) -> bool:
        """
        Определение является ли контент HTML страницей (новое в v1.10.3).

        Args:
            content_type: MIME тип из заголовков
            url: URL для анализа расширения как fallback

        Returns:
            bool: True если это HTML страница
        """
        # Приоритет: Content-Type
        if "text/html" in content_type or "application/xhtml" in content_type:
            return True

        # Проверяем специфические случаи
        if "text/plain" in content_type:
            # Для text/plain проверяем расширение URL
            extension = get_file_extension(url.split("?")[0])  # strip query params
            return extension in ["html", "htm"]

        # Если Content-Type неопределенный или отсутствует
        if not content_type or "application/octet-stream" in content_type:
            extension = get_file_extension(url.split("?")[0])  # strip query params
            return (
                extension in ["html", "htm"] or extension is None
            )  # None означает вероятно динамическая страница

        return False

    def _download_and_extract_file(
        self,
        url: str,
        user_agent: Optional[str] = None,
        extraction_options: Optional[Any] = None,
    ) -> List[Dict[str, Any]]:
        """
        Скачивание файла по URL и его обработка как обычного файла (новое в v1.10.3).

        Args:
            url: URL файла для скачивания
            user_agent: Пользовательский User-Agent
            extraction_options: Настройки извлечения

        Returns:
            List[Dict[str, Any]]: Результат извлечения текста как от /v1/extract/file
        """
        if not requests:
            raise ValueError("requests library not available for file download")

        # Определяем настройки с учетом переданных параметров или значений по умолчанию
        download_timeout = (
            extraction_options.web_page_timeout
            if extraction_options and extraction_options.web_page_timeout is not None
            else settings.FILE_DOWNLOAD_TIMEOUT
        )

        follow_redirects = (
            extraction_options.follow_redirects
            if extraction_options and extraction_options.follow_redirects is not None
            else True
        )

        # Установка User-Agent и заголовков
        headers = {
            "User-Agent": user_agent or settings.DEFAULT_USER_AGENT,
            "Accept": "*/*",
            "Connection": "keep-alive",
        }

        # Создание сессии
        session = requests.Session()
        session.headers.update(headers)

        temp_file_path = None
        try:
            logger.info(f"Скачиваю файл с URL: {url}")

            response = session.get(
                url,
                timeout=download_timeout,
                allow_redirects=follow_redirects,
                stream=True,
            )
            response.raise_for_status()

            # Проверяем размер файла
            content_length = response.headers.get("content-length")
            if content_length and int(content_length) > settings.MAX_FILE_SIZE:
                raise ValueError(
                    f"File too large: {content_length} bytes (max {settings.MAX_FILE_SIZE} bytes)"
                )

            # Определяем имя файла
            filename = self._extract_filename_from_response(response, url)

            # Создаем временный файл
            import tempfile

            suffix = (
                f".{get_file_extension(filename)}"
                if get_file_extension(filename)
                else ""
            )
            with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as temp_file:
                temp_file_path = temp_file.name

                # Скачиваем файл порциями с проверкой размера
                downloaded_size = 0
                for chunk in response.iter_content(chunk_size=8192):
                    if chunk:
                        downloaded_size += len(chunk)
                        if downloaded_size > settings.MAX_FILE_SIZE:
                            raise ValueError(
                                f"File too large: exceeded {settings.MAX_FILE_SIZE} bytes during download"
                            )
                        temp_file.write(chunk)

            logger.info(f"Файл скачан ({downloaded_size} байт): {filename}")

            # Читаем скачанный файл и обрабатываем его как обычный файл
            with open(temp_file_path, "rb") as f:
                file_content = f.read()

            # Используем существующую логику извлечения текста
            return self.extract_text(file_content, filename)

        except Exception as e:
            logger.error(f"Ошибка при скачивании файла {url}: {str(e)}")
            raise ValueError(f"Error downloading file: {str(e)}")
        finally:
            # Очистка временного файла
            if temp_file_path and os.path.exists(temp_file_path):
                try:
                    os.unlink(temp_file_path)
                    logger.debug(f"Временный файл удален: {temp_file_path}")
                except OSError as e:
                    logger.warning(
                        f"Не удалось удалить временный файл {temp_file_path}: {str(e)}"
                    )
            session.close()

    def _extract_filename_from_response(self, response, url: str) -> str:
        """
        Извлечение имени файла из HTTP ответа (новое в v1.10.3).

        Args:
            response: HTTP ответ requests
            url: Исходный URL

        Returns:
            str: Имя файла
        """
        # Пытаемся получить имя файла из заголовка Content-Disposition
        content_disposition = response.headers.get("content-disposition", "")
        if "filename=" in content_disposition:
            import re

            filename_match = re.search(
                r'filename=["\']*([^"\';\r\n]*)', content_disposition
            )
            if filename_match:
                filename = filename_match.group(1).strip()
                if filename:
                    from .utils import sanitize_filename
                    return sanitize_filename(filename)

        # Use last URL segment as filename
        from urllib.parse import unquote, urlparse

        parsed_url = urlparse(url)
        filename = unquote(parsed_url.path.split("/")[-1])

        if not get_file_extension(filename):
            content_type = response.headers.get("content-type", "").lower()
            extension = self._get_extension_from_content_type(content_type)
            if extension:
                filename = f"{filename}.{extension}"

        from .utils import sanitize_filename
        return sanitize_filename(filename) if filename else "downloaded_file"

    def _get_extension_from_content_type(self, content_type: str) -> Optional[str]:
        """
        Определение расширения файла по Content-Type (новое в v1.10.3).

        Args:
            content_type: MIME тип

        Returns:
            Optional[str]: Расширение файла или None
        """
        # Маппинг популярных MIME типов на расширения
        mime_to_extension = settings.MIME_TO_EXTENSION

        # Убираем параметры из Content-Type (например, charset)
        clean_content_type = content_type.split(";")[0].strip()

        return mime_to_extension.get(clean_content_type)

    def extract_from_url(
        self,
        url: str,
        user_agent: Optional[str] = None,
        extraction_options: Optional[Any] = None,
    ) -> List[Dict[str, Any]]:
        """Извлечение текста с веб-страницы или файла по URL (обновлено в v1.10.3)."""
        # Проверка безопасности URL
        if not self._is_safe_url(url):
            raise ValueError(
                "Access to internal IP addresses is prohibited for security reasons"
            )

        try:
            # Шаг 1: Определяем тип контента через HEAD запрос
            content_type, final_url = self._determine_content_type(
                url, user_agent, extraction_options
            )

            # Шаг 2: Выбираем стратегию обработки
            if self._is_html_content(content_type, final_url):
                logger.info(
                    f"URL {final_url} определен как HTML страница (Content-Type: {content_type}), используем веб-экстрактор"
                )
                return self._extract_html_page(
                    final_url, user_agent, extraction_options
                )
            else:
                logger.info(
                    f"URL {final_url} определен как файл (Content-Type: {content_type}), скачиваем и обрабатываем"
                )
                return self._download_and_extract_file(
                    final_url, user_agent, extraction_options
                )

        except Exception as e:
            logger.error(f"Ошибка при обработке URL {url}: {str(e)}")
            raise ValueError(f"Error processing URL: {str(e)}")

    def _extract_html_page(
        self,
        url: str,
        user_agent: Optional[str] = None,
        extraction_options: Optional[Any] = None,
    ) -> List[Dict[str, Any]]:
        """
        Извлечение текста с HTML страницы (выделено из extract_from_url в v1.10.3).

        Args:
            url: URL HTML страницы
            user_agent: Пользовательский User-Agent
            extraction_options: Настройки извлечения

        Returns:
            List[Dict[str, Any]]: Результат извлечения текста со страницы
        """
        html_content = ""
        final_url = url

        # Определяем настройки с учетом переданных параметров или значений по умолчанию
        enable_javascript = (
            extraction_options.enable_javascript
            if extraction_options and extraction_options.enable_javascript is not None
            else settings.ENABLE_JAVASCRIPT
        )

        # Выбираем метод загрузки в зависимости от настроек JavaScript
        if enable_javascript and sync_playwright:
            logger.info("Использую Playwright для загрузки страницы с JS")
            try:
                html_content, final_url = self._extract_page_with_playwright(
                    url, user_agent, extraction_options
                )
            except Exception as e:
                logger.warning(f"Ошибка Playwright: {str(e)}, переключаюсь на requests")
                # Fallback на requests при ошибке Playwright
                html_content, final_url = self._extract_page_with_requests(
                    url, user_agent, extraction_options
                )
        else:
            if enable_javascript and not sync_playwright:
                logger.warning(
                    "JavaScript включен, но Playwright не установлен, использую requests"
                )
            logger.info("Использую requests для загрузки страницы")
            html_content, final_url = self._extract_page_with_requests(
                url, user_agent, extraction_options
            )

        try:
            # Извлечение текста из HTML
            page_text = self._extract_text_from_html(html_content)

            # Поиск и обработка изображений
            image_texts = self._extract_images_from_html(
                html_content, final_url, extraction_options
            )

            # Формирование результата
            results = []

            # Добавляем основной контент страницы
            results.append(
                {
                    "filename": "page_content",
                    "path": final_url,
                    "size": len(html_content.encode("utf-8")),
                    "type": "html",
                    "text": page_text,
                }
            )

            # Добавляем тексты с изображений
            results.extend(image_texts)

            return results

        except Exception as e:
            raise ValueError(f"Error processing HTML page: {str(e)}")

    def _extract_page_with_requests(
        self,
        url: str,
        user_agent: Optional[str] = None,
        extraction_options: Optional[Any] = None,
    ) -> tuple[str, str]:
        """
        Извлечение HTML контента страницы с помощью requests (без JS, обновлено в v1.10.2).

        Args:
            url: URL страницы
            user_agent: Пользовательский User-Agent
            extraction_options: Настройки извлечения

        Returns:
            tuple[str, str]: (html_content, final_url)
        """
        if not requests:
            raise ValueError("requests library not available for web extraction")

        # Определяем настройки с учетом переданных параметров или значений по умолчанию
        web_page_timeout = (
            extraction_options.web_page_timeout
            if extraction_options and extraction_options.web_page_timeout is not None
            else settings.WEB_PAGE_TIMEOUT
        )

        follow_redirects = (
            extraction_options.follow_redirects
            if extraction_options and extraction_options.follow_redirects is not None
            else True
        )

        # Установка User-Agent
        headers = {
            "User-Agent": user_agent or settings.DEFAULT_USER_AGENT,
            "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8",
            "Accept-Language": "ru,en;q=0.5",
            "Accept-Encoding": "gzip, deflate",
            "Connection": "keep-alive",
        }

        try:
            # Загрузка страницы с таймаутом
            response = requests.get(
                url,
                headers=headers,
                timeout=web_page_timeout,
                allow_redirects=follow_redirects,
                stream=False,
            )
            response.raise_for_status()

            # Автоопределение кодировки
            response.encoding = response.apparent_encoding or "utf-8"
            html_content = response.text
            final_url = response.url

            logger.info(
                f"HTML получен через requests, размер: {len(html_content)} символов"
            )
            return html_content, final_url

        except requests.RequestException as e:
            if "timeout" in str(e).lower():
                raise ValueError(f"Page loading timeout: {str(e)}")
            elif "connection" in str(e).lower():
                raise ValueError(f"Connection error: {str(e)}")
            else:
                raise ValueError(f"Failed to load page: {str(e)}")

    def _is_safe_url(self, url: str) -> bool:
        """Проверка безопасности URL (защита от SSRF)."""
        try:
            parsed_url = urlparse(url)

            # Проверяем схему URL
            if not self._check_url_scheme(parsed_url.scheme):
                return False

            hostname = parsed_url.hostname
            if not hostname:
                logger.warning(f"No hostname in URL: {url}")
                return False

            # Проверяем заблокированные хосты
            if not self._check_hostname_not_blocked(hostname, url):
                return False

            # Получаем IP-адреса хоста
            ips = self._resolve_hostname_ips(hostname)
            if not ips:
                return False

            # Проверяем безопасность всех IP-адресов
            return self._check_all_ips_safe(ips, url)

        except Exception as e:
            logger.warning(f"Error checking URL safety: {str(e)}")
            # Fail-closed: в случае ошибки блокируем доступ
            return False

    def _check_url_scheme(self, scheme: str) -> bool:
        """Проверка схемы URL."""
        if scheme not in ["http", "https"]:
            logger.warning(f"Unsupported URL scheme: {scheme}")
            return False
        return True

    def _check_hostname_not_blocked(self, hostname: str, url: str) -> bool:
        """Проверка, что hostname не заблокирован."""
        blocked_hostnames = settings.BLOCKED_HOSTNAMES.split(",")
        hostname_lower = hostname.lower()

        for blocked_hostname in blocked_hostnames:
            blocked_hostname = blocked_hostname.strip().lower()
            if blocked_hostname and hostname_lower == blocked_hostname:
                logger.warning(f"Blocked hostname {hostname} for URL {url}")
                return False
        return True

    def _resolve_hostname_ips(self, hostname: str) -> list:
        """Разрешение IP-адресов для hostname."""
        import socket

        try:
            addr_info = socket.getaddrinfo(
                hostname, None, socket.AF_UNSPEC, socket.SOCK_STREAM
            )
            return [info[4][0] for info in addr_info]
        except socket.gaierror as e:
            logger.warning(f"DNS resolution failed for {hostname}: {str(e)}")
            return []

    def _check_all_ips_safe(self, ips: list, url: str) -> bool:
        """Проверка безопасности всех IP-адресов."""
        for ip_str in ips:
            if not self._check_single_ip_safe(ip_str, url):
                return False
        return True

    def _check_single_ip_safe(self, ip_str: str, url: str) -> bool:
        """Проверка безопасности одного IP-адреса."""
        try:
            ip_obj = ipaddress.ip_address(ip_str)

            # Проверяем специальные адреса
            if self._is_special_ip_unsafe(ip_obj, ip_str, url):
                return False

            # Проверяем заблокированные диапазоны
            if self._is_ip_in_blocked_ranges(ip_obj, ip_str, url):
                return False

            # Проверяем metadata service
            if self._is_metadata_service_ip(ip_obj, ip_str, url):
                return False

            # Проверяем Docker bridge
            if self._is_docker_bridge_ip(ip_obj, ip_str, url):
                return False

            return True

        except ValueError as e:
            logger.warning(f"Invalid IP address {ip_str}: {str(e)}")
            return True  # Невалидный IP не блокируем

    def _is_special_ip_unsafe(self, ip_obj, ip_str: str, url: str) -> bool:
        """Проверка на специальные небезопасные IP."""
        if ip_obj.is_loopback or ip_obj.is_private or ip_obj.is_link_local:
            logger.warning(
                f"Blocked special IP {ip_str} (loopback/private/link-local) for URL {url}"
            )
            return True
        return False

    def _is_ip_in_blocked_ranges(self, ip_obj, ip_str: str, url: str) -> bool:
        """Проверка IP на принадлежность заблокированным диапазонам."""
        blocked_ranges = settings.BLOCKED_IP_RANGES.split(",")

        for range_str in blocked_ranges:
            range_str = range_str.strip()
            if not range_str:
                continue

            try:
                network = ipaddress.ip_network(range_str, strict=False)
                if ip_obj in network:
                    logger.warning(
                        f"Blocked IP {ip_str} in range {range_str} for URL {url}"
                    )
                    return True
            except ValueError:
                continue
        return False

    def _is_metadata_service_ip(self, ip_obj, ip_str: str, url: str) -> bool:
        """Проверка на metadata service IP."""
        if str(ip_obj) == "169.254.169.254":
            logger.warning(f"Blocked metadata service IP {ip_str} for URL {url}")
            return True
        return False

    def _is_docker_bridge_ip(self, ip_obj, ip_str: str, url: str) -> bool:
        """Проверка на Docker bridge gateway IP."""
        if ip_obj.version == 4:
            octets = str(ip_obj).split(".")
            if (
                octets[0] == "172"
                and 16 <= int(octets[1]) <= 31
                and octets[2] == "0"
                and octets[3] == "1"
            ):
                logger.warning(f"Blocked Docker bridge gateway {ip_str} for URL {url}")
                return True
        return False

    def _extract_text_from_html(self, html_content: str) -> str:
        """Извлечение текста из HTML контента."""
        if not BeautifulSoup:
            raise ValueError("BeautifulSoup not available for HTML parsing")

        try:
            soup = BeautifulSoup(html_content, "lxml")

            # Удаляем скрипты, стили и другие нетекстовые элементы
            for script in soup(["script", "style", "nav", "header", "footer", "aside"]):
                script.decompose()

            # Извлекаем текст
            text = soup.get_text()

            # Очистка текста
            lines = []
            for line in text.splitlines():
                line = line.strip()
                if line:
                    lines.append(line)

            return "\n".join(lines)

        except Exception as e:
            logger.error(f"Error extracting text from HTML: {str(e)}")
            raise ValueError(f"HTML parsing error: {str(e)}")

    def _extract_images_from_html(
        self, html_content: str, base_url: str, extraction_options: Optional[Any] = None
    ) -> List[Dict[str, Any]]:
        """Извлечение и обработка изображений со страницы (обновлено в v1.10.2)."""
        if not BeautifulSoup or not Image:
            return []

        # Настройка параметров извлечения
        options = self._setup_image_extraction_options(extraction_options)
        if not options["process_images"]:
            logger.info("Обработка изображений отключена в настройках извлечения")
            return []

        try:
            # Парсинг изображений из HTML
            img_tags = self._parse_images_from_html(
                html_content, options["max_images_per_page"]
            )
            if not img_tags:
                return []

            # Категоризация изображений
            base64_images, url_images = self._categorize_images(
                img_tags, options["enable_base64_images"]
            )
            logger.info(
                f"Найдено изображений: {len(url_images)} URL, {len(base64_images)} base64"
            )

            results = []
            # Обработка изображений
            results.extend(
                self._process_base64_images(base64_images, extraction_options)
            )
            results.extend(
                self._process_url_images(url_images, base_url, extraction_options)
            )

            return results

        except Exception as e:
            logger.warning(f"Error extracting images from HTML: {str(e)}")
            return []

    def _setup_image_extraction_options(
        self, extraction_options: Optional[Any]
    ) -> dict:
        """Настройка параметров извлечения изображений."""
        return {
            "process_images": (
                extraction_options.process_images
                if extraction_options and extraction_options.process_images is not None
                else True
            ),
            "max_images_per_page": (
                extraction_options.max_images_per_page
                if extraction_options
                and extraction_options.max_images_per_page is not None
                else settings.MAX_IMAGES_PER_PAGE
            ),
            "enable_base64_images": (
                extraction_options.enable_base64_images
                if extraction_options
                and extraction_options.enable_base64_images is not None
                else settings.ENABLE_BASE64_IMAGES
            ),
        }

    def _parse_images_from_html(self, html_content: str, max_images: int) -> list:
        """Парсинг изображений из HTML контента."""
        soup = BeautifulSoup(html_content, "lxml")
        img_tags = soup.find_all("img", src=True)
        return img_tags[:max_images]

    def _categorize_images(self, img_tags: list, enable_base64: bool) -> tuple:
        """Категоризация изображений на base64 и URL."""
        base64_images = []
        url_images = []

        for img_tag in img_tags:
            img_src = img_tag.get("src", "")
            if img_src.startswith("data:image/") and enable_base64:
                base64_images.append(img_tag)
            else:
                url_images.append(img_tag)

        return base64_images, url_images

    def _process_base64_images(
        self, base64_images: list, extraction_options: Optional[Any]
    ) -> list:
        """Обработка base64 изображений."""
        results = []
        for img_tag in base64_images:
            try:
                result = self._process_base64_image(img_tag, extraction_options)
                if result:
                    results.append(result)
            except Exception as e:
                logger.warning(f"Error processing base64 image: {str(e)}")
        return results

    def _process_url_images(
        self, url_images: list, base_url: str, extraction_options: Optional[Any]
    ) -> list:
        """Обработка URL изображений."""
        if not url_images:
            return []

        results = []
        image_download_timeout = (
            extraction_options.image_download_timeout
            if extraction_options
            and extraction_options.image_download_timeout is not None
            else settings.IMAGE_DOWNLOAD_TIMEOUT
        )

        # Обработка изображений группами по 2
        for i in range(0, len(url_images), 2):
            batch = url_images[i : i + 2]
            batch_results = self._process_images_batch(
                batch, base_url, extraction_options, image_download_timeout
            )
            results.extend(batch_results)

        return results

    def _process_images_batch(
        self,
        batch: list,
        base_url: str,
        extraction_options: Optional[Any],
        timeout: int,
    ) -> list:
        """Обработка группы изображений параллельно."""
        batch_results = []

        with concurrent.futures.ThreadPoolExecutor(max_workers=2) as executor:
            futures = []
            for img_tag in batch:
                future = executor.submit(
                    self._process_single_image,
                    img_tag,
                    base_url,
                    extraction_options,
                )
                futures.append(future)

            for future in concurrent.futures.as_completed(futures):
                try:
                    result = future.result(timeout=timeout + 5)
                    if result:
                        batch_results.append(result)
                except Exception as e:
                    logger.warning(f"Error processing image: {str(e)}")

        return batch_results

    def _process_single_image(
        self, img_tag, base_url: str, extraction_options: Optional[Any] = None
    ) -> Optional[Dict[str, Any]]:
        """Обработка одного изображения (обновлено в v1.10.2)."""
        try:
            img_src = img_tag.get("src", "")
            logger.info(f"Processing image: {img_src}")
            if not img_src:
                logger.warning("Image has no src attribute")
                return None

            # Преобразуем относительный URL в абсолютный
            img_url = urljoin(base_url, img_src)
            logger.info(f"Full image URL: {img_url}")

            # Проверяем безопасность URL изображения
            if not self._is_safe_url(img_url):
                logger.warning(f"Blocked image URL: {img_url}")
                return None

            # Определяем настройки
            image_download_timeout = (
                extraction_options.image_download_timeout
                if extraction_options
                and extraction_options.image_download_timeout is not None
                else settings.IMAGE_DOWNLOAD_TIMEOUT
            )

            min_image_size_for_ocr = (
                extraction_options.min_image_size_for_ocr
                if extraction_options
                and extraction_options.min_image_size_for_ocr is not None
                else settings.MIN_IMAGE_SIZE_FOR_OCR
            )

            # Загрузка изображения
            headers = {"User-Agent": settings.DEFAULT_USER_AGENT, "Referer": base_url}

            response = requests.get(
                img_url, headers=headers, timeout=image_download_timeout, stream=True
            )
            response.raise_for_status()

            # Проверяем размер изображения
            img_content = response.content
            logger.info(f"Image content size: {len(img_content)} bytes")
            if len(img_content) == 0:
                logger.warning("Image content is empty")
                return None

            # Открываем изображение для проверки размеров
            with Image.open(io.BytesIO(img_content)) as img:
                width, height = img.size
                logger.info(
                    f"Image dimensions: {width}x{height} = {width * height} pixels (min required: {settings.MIN_IMAGE_SIZE_FOR_OCR})"
                )

                # Проверяем минимальный размер
                if width * height < min_image_size_for_ocr:
                    logger.warning(
                        f"Image too small for OCR: {width * height} < {min_image_size_for_ocr}"
                    )
                    return None

                # OCR изображения
                logger.info(f"Starting OCR for image: {img_url}")
                text = self._safe_tesseract_ocr(img)
                logger.info(f"OCR result length: {len(text) if text else 0} characters")

                if not text or not text.strip():
                    logger.warning("No text found in image")
                    return None

                # Извлекаем имя файла из URL
                from .utils import get_extension_from_mime

                filename = os.path.basename(urlparse(img_url).path) or "image"
                if "." not in filename:
                    # Определяем расширение по MIME-типу через утилиту
                    content_type = response.headers.get("content-type", "").lower()
                    extension = get_extension_from_mime(
                        content_type, settings.SUPPORTED_FORMATS
                    )

                    if extension:
                        filename += f".{extension}"
                    else:
                        # Если MIME-тип не поддерживается, игнорируем изображение
                        logger.warning(f"Unsupported image MIME type: {content_type}")
                        return None

                return {
                    "filename": filename,
                    "path": img_url,
                    "size": len(img_content),
                    "type": filename.split(".")[-1].lower(),
                    "text": text.strip(),
                }

        except Exception as e:
            logger.warning(f"Error processing image {img_tag.get('src', '')}: {str(e)}")
            return None

    def _process_base64_image(
        self, img_tag, extraction_options: Optional[Any] = None
    ) -> Optional[Dict[str, Any]]:
        """Обработка base64 изображения из data URI (обновлено в v1.10.2)."""
        try:
            from .utils import (
                decode_base64_image,
                extract_mime_from_base64_data_uri,
                get_extension_from_mime,
            )

            img_src = img_tag.get("src", "")
            logger.info(f"Processing base64 image: {img_src[:50]}...")

            if not img_src.startswith("data:image/"):
                logger.warning("Invalid base64 image format")
                return None

            # Определяем настройки
            min_image_size_for_ocr = (
                extraction_options.min_image_size_for_ocr
                if extraction_options
                and extraction_options.min_image_size_for_ocr is not None
                else settings.MIN_IMAGE_SIZE_FOR_OCR
            )

            # Извлекаем MIME-тип
            mime_type = extract_mime_from_base64_data_uri(img_src)
            if not mime_type:
                logger.warning("Could not extract MIME type from base64 image")
                return None

            # Определяем расширение файла
            extension = get_extension_from_mime(mime_type, settings.SUPPORTED_FORMATS)
            if not extension:
                logger.warning(f"Unsupported image MIME type: {mime_type}")
                return None

            # Декодируем base64 изображение
            img_content = decode_base64_image(img_src)
            if not img_content:
                logger.warning("Failed to decode base64 image")
                return None

            logger.info(f"Base64 image decoded, size: {len(img_content)} bytes")

            # Открываем изображение для проверки размеров
            with Image.open(io.BytesIO(img_content)) as img:
                width, height = img.size
                logger.info(
                    f"Base64 image dimensions: {width}x{height} = {width * height} pixels (min required: {min_image_size_for_ocr})"
                )

                # Проверяем минимальный размер
                if width * height < min_image_size_for_ocr:
                    logger.warning(
                        f"Base64 image too small for OCR: {width * height} < {min_image_size_for_ocr}"
                    )
                    return None

                # OCR изображения
                logger.info("Starting OCR for base64 image")
                text = self._safe_tesseract_ocr(img)
                logger.info(f"OCR result length: {len(text) if text else 0} characters")

                if not text or not text.strip():
                    logger.warning("No text found in base64 image")
                    return None

                # Формируем имя файла
                filename = f"base64_image.{extension}"

                return {
                    "filename": filename,
                    "path": f"data:image/{extension};base64,[base64_data]",
                    "size": len(img_content),
                    "type": extension,
                    "text": text.strip(),
                }

        except Exception as e:
            logger.warning(f"Error processing base64 image: {str(e)}")
            return None
